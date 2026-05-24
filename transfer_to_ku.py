"""Transfer PMG_presentation.pptx content into the KU-format template.

Produces PMG_ku_final.pptx by:
1. Opening PMG_ku_presentation.pptx as the base (preserves KU master/theme).
2. Stripping the 9 example slides via XML manipulation.
3. Adding 11 new slides mapped from the source, using appropriate KU layouts.
"""

from __future__ import annotations

import copy
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
_HERE = Path(__file__).parent
SRC_PATH = _HERE / "PMG_presentation.pptx"
KU_PATH = _HERE / "PMG_ku_presentation.pptx"
OUT_PATH = _HERE / "PMG_ku_final.pptx"

# KU layout indices (confirmed from inspection)
_LAYOUT_TITLE_LARGE = 1      # Title seal large
_LAYOUT_BULLET = 14          # Title and one bullet list
_LAYOUT_EMPTY = 11           # Empty
_LAYOUT_TWO_IMAGES = 19      # Two images
_LAYOUT_TITLE_IMAGE = 9      # Title and image
_LAYOUT_QUESTIONS = 28       # Questions and comments

# Footer/date/slide-number placeholder indices — never copied from source
_FOOTER_IDXS = {10, 11, 12, 13, 16, 17, 18}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _clear_slide_list(prs: Presentation) -> None:
    """Remove all existing slides from the presentation.

    Drops both the slide-ID XML entries and the OPC part relationships so
    that python-pptx does not encounter duplicate part-name collisions when
    adding new slides.
    """
    # Collect rIds for all slide relationships before mutating anything
    slide_rids = [
        rId
        for rId, rel in prs.part.rels.items()
        if "slide" in rel.reltype and "slideLayout" not in rel.reltype
        and "slideMaster" not in rel.reltype
    ]
    # Remove from the XML slide-ID list
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        xml_slides.remove(sld)
    # Drop the corresponding OPC relationships so part names are freed
    for rId in slide_rids:
        prs.part.rels.pop(rId)


def _add_slide(prs: Presentation, layout_idx: int):
    """Add a blank slide with the given layout index and return it."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def _get_placeholder(slide, idx: int):
    """Return placeholder by idx, or None if absent."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _set_title(slide, text: str) -> None:
    """Write *text* into the title placeholder (idx=0) if present."""
    ph = _get_placeholder(slide, 0)
    if ph is None or not text.strip():
        return
    tf = ph.text_frame
    tf.clear()
    tf.paragraphs[0].text = text.rstrip("\x0b").strip()


def _set_body_bullets(slide, src_shape) -> None:
    """Copy bullet paragraphs from *src_shape* into the body placeholder (idx=1).

    Preserves paragraph indent level. Each run's text is concatenated so that
    multi-run paragraphs land in a single paragraph on the destination.
    """
    ph = _get_placeholder(slide, 1)
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()
    first = True
    for src_para in src_shape.text_frame.paragraphs:
        text = src_para.text
        if not text.strip() and not first:
            continue  # skip trailing empty paragraphs
        if first:
            dst_para = tf.paragraphs[0]
            first = False
        else:
            dst_para = tf.add_paragraph()
        dst_para.text = text
        dst_para.level = src_para.level


def _extract_png(shape) -> bytes | None:
    """Return raw PNG bytes from a PICTURE shape, or None for Grafik shapes."""
    try:
        return shape.image.blob
    except (ValueError, AttributeError):
        return None


def _write_tmp_image(blob: bytes, idx: int) -> str:
    """Write image blob to /tmp and return path."""
    path = f"/tmp/slide_img_{idx}.png"
    with open(path, "wb") as f:
        f.write(blob)
    return path


def _add_picture_from_shape(slide, shape, img_idx: int) -> bool:
    """Add a picture to *slide* at the same geometry as *shape*.

    Returns True on success, False if the shape has no embedded image.
    """
    blob = _extract_png(shape)
    if blob is None:
        return False
    tmp = _write_tmp_image(blob, img_idx)
    slide.shapes.add_picture(
        tmp, shape.left, shape.top, shape.width, shape.height
    )
    return True


def _copy_shape_xml(dst_slide, shape) -> None:
    """Deep-copy a shape's XML element into *dst_slide*'s spTree.

    Grafik/EMF shapes that have no embedded blob are handled this way.
    Prints a warning if the copy raises.
    """
    try:
        sp_tree = dst_slide.shapes._spTree
        sp_tree.append(copy.deepcopy(shape._element))
    except Exception as exc:
        print(f"  WARNING: XML copy of '{shape.name}' failed: {exc}", file=sys.stderr)


def _add_textbox_bold(slide, shape) -> None:
    """Recreate a text-box shape with bold 14pt text at the same geometry."""
    txBox = slide.shapes.add_textbox(
        shape.left, shape.top, shape.width, shape.height
    )
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = shape.text_frame.text
    run = p.runs[0] if p.runs else p.add_run()
    run.text = shape.text_frame.text
    run.font.bold = True
    run.font.size = Pt(14)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _build_pause_slide(prs: Presentation) -> None:
    """Slide 1 — blank pause/cover slide."""
    _add_slide(prs, _LAYOUT_TITLE_LARGE)


def _build_title_slide(prs: Presentation, src_slide) -> None:
    """Slide 2 — main title + subtitle."""
    new_slide = _add_slide(prs, _LAYOUT_TITLE_LARGE)
    # idx=0 on Title seal large is the main title
    title_ph = _get_placeholder(new_slide, 0)
    subtitle_ph = _get_placeholder(new_slide, 1)

    for shape in src_slide.shapes:
        try:
            ph_idx = shape.placeholder_format.idx
        except ValueError:
            continue
        if ph_idx == 0 and title_ph is not None:
            title_ph.text_frame.clear()
            title_ph.text_frame.paragraphs[0].text = (
                shape.text_frame.text.rstrip("\x0b").strip()
            )
        elif ph_idx == 1 and subtitle_ph is not None:
            subtitle_ph.text_frame.clear()
            first = True
            for src_para in shape.text_frame.paragraphs:
                if first:
                    dst_para = subtitle_ph.text_frame.paragraphs[0]
                    first = False
                else:
                    dst_para = subtitle_ph.text_frame.add_paragraph()
                dst_para.text = src_para.text


def _build_bullet_slide(
    prs: Presentation, src_slide, img_idx_start: int = 0
) -> int:
    """Slides 3, 4, 6 — title + bullet list + optional PNG image.

    Returns the next available image temp-file index.
    """
    new_slide = _add_slide(prs, _LAYOUT_BULLET)
    img_idx = img_idx_start

    title_text = ""
    body_shape = None
    image_shapes = []

    for shape in src_slide.shapes:
        try:
            ph_idx = shape.placeholder_format.idx
        except ValueError:
            # Non-placeholder shape (e.g. picture)
            if shape.shape_type == 13:
                image_shapes.append(shape)
            continue
        if ph_idx in _FOOTER_IDXS:
            continue
        if ph_idx == 0:
            title_text = shape.text_frame.text.rstrip("\x0b").strip()
        elif ph_idx == 1:
            body_shape = shape

    _set_title(new_slide, title_text)
    if body_shape is not None:
        _set_body_bullets(new_slide, body_shape)

    for img_shape in image_shapes:
        success = _add_picture_from_shape(new_slide, img_shape, img_idx)
        if not success:
            _copy_shape_xml(new_slide, img_shape)
        img_idx += 1

    return img_idx


def _build_empty_slide(prs: Presentation, src_slide) -> None:
    """Slides 8, 9 — empty layout; copy Grafik shapes via XML."""
    new_slide = _add_slide(prs, _LAYOUT_EMPTY)
    for shape in src_slide.shapes:
        try:
            shape.placeholder_format  # raises if not a placeholder
            continue  # skip all placeholder shapes
        except ValueError:
            pass
        if shape.shape_type == 13:
            blob = _extract_png(shape)
            if blob is not None:
                tmp = _write_tmp_image(blob, id(shape))
                new_slide.shapes.add_picture(
                    tmp, shape.left, shape.top, shape.width, shape.height
                )
            else:
                _copy_shape_xml(new_slide, shape)


def _build_two_image_slide(
    prs: Presentation, src_slide, img_idx_start: int
) -> int:
    """Slide 7 — title, two PNG images, two bold label text boxes."""
    new_slide = _add_slide(prs, _LAYOUT_TWO_IMAGES)
    img_idx = img_idx_start

    title_text = ""
    image_shapes = []
    textbox_shapes = []

    for shape in src_slide.shapes:
        try:
            ph_idx = shape.placeholder_format.idx
        except ValueError:
            if shape.shape_type == 13:
                image_shapes.append(shape)
            elif shape.has_text_frame:
                textbox_shapes.append(shape)
            continue
        if ph_idx == 0:
            title_text = shape.text_frame.text.rstrip("\x0b").strip()

    _set_title(new_slide, title_text)

    for img_shape in image_shapes:
        success = _add_picture_from_shape(new_slide, img_shape, img_idx)
        if not success:
            _copy_shape_xml(new_slide, img_shape)
        img_idx += 1

    for tb_shape in textbox_shapes:
        _add_textbox_bold(new_slide, tb_shape)

    return img_idx


def _build_title_image_slide(
    prs: Presentation, src_slide, img_idx_start: int
) -> int:
    """Slide 10 — title + single image."""
    new_slide = _add_slide(prs, _LAYOUT_TITLE_IMAGE)
    img_idx = img_idx_start

    title_text = ""
    image_shapes = []

    for shape in src_slide.shapes:
        try:
            ph_idx = shape.placeholder_format.idx
        except ValueError:
            if shape.shape_type == 13:
                image_shapes.append(shape)
            continue
        if ph_idx == 0:
            title_text = shape.text_frame.text.rstrip("\x0b").strip()

    _set_title(new_slide, title_text)

    for img_shape in image_shapes:
        success = _add_picture_from_shape(new_slide, img_shape, img_idx)
        if not success:
            _copy_shape_xml(new_slide, img_shape)
        img_idx += 1

    return img_idx


def _build_questions_slide(prs: Presentation) -> None:
    """Slide 11 — KU 'Questions and comments' closing slide."""
    _add_slide(prs, _LAYOUT_QUESTIONS)


# ---------------------------------------------------------------------------
# Main transfer routine
# ---------------------------------------------------------------------------

def transfer() -> int:
    """Execute the full transfer and return the final slide count."""
    src = Presentation(str(SRC_PATH))
    prs = Presentation(str(KU_PATH))

    _clear_slide_list(prs)

    src_slides = src.slides
    img_idx = 0  # monotonic counter for /tmp image filenames

    # Slide 1 — pause/cover (blank)
    _build_pause_slide(prs)

    # Slide 2 — main title
    _build_title_slide(prs, src_slides[1])

    # Slide 3 — "What did Guha et al. do?" with inline image
    img_idx = _build_bullet_slide(prs, src_slides[2], img_idx)

    # Slide 4 — "Three Methodological Errors"
    img_idx = _build_bullet_slide(prs, src_slides[3], img_idx)

    # Slide 5 (source) is blank — skip

    # Slide 6 — "Replication — What We Did"
    img_idx = _build_bullet_slide(prs, src_slides[5], img_idx)

    # Slide 7 — two images + bold labels
    img_idx = _build_two_image_slide(prs, src_slides[6], img_idx)

    # Slide 8 — Grafik shapes (XML copy)
    _build_empty_slide(prs, src_slides[7])

    # Slide 9 — Grafik shape (XML copy)
    _build_empty_slide(prs, src_slides[8])

    # Slide 10 — single PNG
    img_idx = _build_title_image_slide(prs, src_slides[9], img_idx)

    # Slide 11 — questions/comments closing slide
    _build_questions_slide(prs)

    prs.save(str(OUT_PATH))
    return len(prs.slides)


if __name__ == "__main__":
    n_slides = transfer()
    print(f"Saved {OUT_PATH} with {n_slides} slides.")
